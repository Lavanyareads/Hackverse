"""
Granite Guardian - Validation Layer (Task 4)
----------------------------------------------
Requires Ollama running locally with:
    ollama pull granite4.1:8b

Run the full test suite (9 cases, for development/regression testing):
    python granite_guardian.py

Run just the curated demo subset (for live presentation):
    python granite_guardian.py --demo

Run against real pipeline files (Task 3's output + Task 2's cleaned input):
    python granite_guardian.py --files
    python granite_guardian.py --files path/to/generated_result.json path/to/cleaned_input.json
"""

import requests
import json
import sys
import os
from datetime import datetime

MODEL = "granite4.1:8b"  # text judge - change to match whatever tag you pulled

# Where every real run (--files mode) writes generated_output.json /
# final.json / retry.json (fixed names, overwritten each run), plus a
# running JSONL log (run_log.jsonl) of every run's outcome. The
# test/demo batch mode (not --files) still writes its own timestamped file.
RESULTS_DIR = "guardian_results"
RUN_LOG_PATH = os.path.join(RESULTS_DIR, "run_log.jsonl")
os.makedirs(RESULTS_DIR, exist_ok=True)


class GuardianConnectionError(Exception):
    """Raised when Guardian can't reach the local Granite model via Ollama."""
    pass


def call_granite(prompt, timeout=180):
    """Send a prompt to the local Granite model via Ollama's API and return the text reply.

    Raises GuardianConnectionError with a clear, actionable message if Ollama
    isn't running, the request times out, or the model isn't available -
    instead of letting a raw connection error crash the whole script.

    timeout defaults to 180s because the FIRST call after Ollama starts (or
    after ~5 min idle) has to load the whole model into memory before it can
    even start generating - on CPU-only laptops this can take a while.
    keep_alive tells Ollama to keep the model loaded for 10 minutes after
    each call, so back-to-back calls (like your test suite) stay fast.
    """
    try:
        response = requests.post(
            "http://localhost:11434/api/chat",
            json={
                "model": MODEL,
                "messages": [{"role": "user", "content": prompt}],
                "stream": False,
                "keep_alive": "10m",
            },
            timeout=timeout,
        )
    except requests.exceptions.ConnectionError:
        raise GuardianConnectionError(
            "Could not connect to Ollama at localhost:11434. Is it running? "
            "Try 'ollama serve' in a separate terminal, then try again."
        )
    except requests.exceptions.Timeout:
        raise GuardianConnectionError(
            "Ollama did not respond within " + str(timeout) + " seconds. "
            "The model may still be loading, or the machine may be under heavy load."
        )
    except requests.exceptions.RequestException as e:
        raise GuardianConnectionError("Unexpected error contacting Ollama: " + str(e))

    if response.status_code == 404:
        raise GuardianConnectionError(
            "Ollama responded, but the model '" + MODEL + "' was not found. "
            "Pull it first with: ollama pull " + MODEL
        )
    if response.status_code != 200:
        raise GuardianConnectionError(
            "Ollama returned an error (status " + str(response.status_code) + "): " + response.text
        )

    try:
        return response.json()["message"]["content"]
    except (KeyError, ValueError) as e:
        raise GuardianConnectionError("Ollama responded in an unexpected format: " + str(e))


def check_ollama_status():
    """Quick pre-flight check you can run before a demo: confirms Ollama is
    reachable and the configured model is actually pulled. Prints a clear
    status message either way and returns True/False."""
    try:
        response = requests.get("http://localhost:11434/api/tags", timeout=5)
        response.raise_for_status()
    except requests.exceptions.RequestException:
        print("Could not reach Ollama at localhost:11434.")
        print("Make sure Ollama is installed and running (try 'ollama serve' in a terminal).")
        return False

    available = [m.get("name", "") for m in response.json().get("models", [])]
    if any(MODEL in name for name in available):
        print("Ollama is running and '" + MODEL + "' is available. Ready to go.")
        text_ok = True
    else:
        print("Ollama is running, but '" + MODEL + "' isn't pulled yet.")
        print("Run: ollama pull " + MODEL)
        text_ok = False

    return text_ok


def warm_up_model():
    """Sends one trivial prompt to force Ollama to load the model into memory
    now, so the first REAL test case isn't the one that pays the slow
    cold-load cost. Call this once, right after check_ollama_status()."""
    print("Warming up " + MODEL + " (loading into memory - this can take a while on the first run)...")
    try:
        call_granite("Reply with only the word: ready", timeout=180)
        print("Model loaded and warm.\n")
        return True
    except GuardianConnectionError as e:
        print("Warm-up failed: " + str(e))
        return False


# ============================================================================
# FILE CONTENT EXTRACTION
# ----------------------------------------------------------------------------
# Turns a real generated file (pptx, docx, pdf, xlsx, etc.) into text Guardian
# can judge with the normal 9 checks. Each extractor imports its library
# lazily, so a missing library only breaks that one format, not the whole
# module.
# ============================================================================

TEXT_EXTRACTABLE_FORMATS = {".pptx", ".docx", ".pdf", ".xlsx", ".xls", ".txt", ".csv", ".md"}


def _extract_pptx_text(file_path):
    from pptx import Presentation
    prs = Presentation(file_path)
    slides = []
    for i, slide in enumerate(prs.slides, 1):
        lines = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = "".join(run.text for run in para.runs)
                    if line.strip():
                        lines.append(line.strip())
        slides.append("Slide " + str(i) + ": " + " | ".join(lines))
    return "\n".join(slides)


def _extract_docx_text(file_path):
    import docx
    doc = docx.Document(file_path)
    parts = [p.text for p in doc.paragraphs if p.text.strip()]
    for table in doc.tables:
        for row in table.rows:
            cells = [c.text.strip() for c in row.cells]
            if any(cells):
                parts.append(" | ".join(cells))
    return "\n".join(parts)


def _extract_pdf_text(file_path):
    from pypdf import PdfReader
    reader = PdfReader(file_path)
    pages = []
    for i, page in enumerate(reader.pages, 1):
        text = (page.extract_text() or "").strip()
        pages.append("Page " + str(i) + ": " + text)
    return "\n".join(pages)


def _extract_xlsx_text(file_path):
    from openpyxl import load_workbook
    wb = load_workbook(file_path, data_only=True)
    lines = []
    for sheet in wb.worksheets:
        lines.append("Sheet: " + sheet.title)
        for row in sheet.iter_rows(values_only=True):
            if any(cell is not None for cell in row):
                lines.append(" | ".join("" if c is None else str(c) for c in row))
    return "\n".join(lines)


def _extract_plain_text(file_path):
    with open(file_path, "r", encoding="utf-8", errors="replace") as f:
        return f.read()


_TEXT_EXTRACTORS = {
    ".pptx": _extract_pptx_text,
    ".docx": _extract_docx_text,
    ".pdf": _extract_pdf_text,
    ".xlsx": _extract_xlsx_text,
    ".xls": _extract_xlsx_text,  # note: openpyxl can't read legacy .xls (pre-2007) - flagged in the error if it fails
    ".txt": _extract_plain_text,
    ".csv": _extract_plain_text,
    ".md": _extract_plain_text,
}


def extract_file_content(file_path):
    """Given a path to a generated file, figure out how to read it.

    Returns a dict:
      {"kind": "text", "content": "...", "error": None}          - text-extractable formats
      {"kind": "unsupported"/"error", "content": None, "error": "..."} - can't process it

    Never raises - always returns something guardian_check_from_orchestrator
    can act on, even for missing files or missing libraries.
    """
    if not os.path.exists(file_path):
        return {"kind": "error", "content": None, "error": "File not found: " + file_path}

    ext = os.path.splitext(file_path)[1].lower()

    if ext in _TEXT_EXTRACTORS:
        try:
            return {"kind": "text", "content": _TEXT_EXTRACTORS[ext](file_path), "error": None}
        except ImportError as e:
            return {"kind": "error", "content": None,
                    "error": "Missing library to read " + ext + " files: " + str(e)}
        except Exception as e:
            return {"kind": "error", "content": None,
                    "error": "Failed to read " + file_path + ": " + str(e)}

    return {"kind": "unsupported", "content": None,
            "error": "No extractor available yet for file type '" + ext + "'"}


GUARDIAN_PROMPT_TEMPLATE = """You are a strict quality reviewer for an AI system. Evaluate the AI OUTPUT against the USER REQUEST, UPLOADED FILES, and SOURCE TEXT using the checks below. Judge each check ONLY on its own specific criterion - do not let an issue in one check affect another check's verdict.

ISOLATION RULE: factual_accuracy is the ONLY check that may fail an output because a number or fact is wrong. For every other check (requirement_match, missing_information, completeness, formatting, relevance), take any stated figures at face value as if they were correct, and judge purely on that check's own criterion below. Example: if the output claims "$6.5M revenue" and that number is wrong, only factual_accuracy should fail - requirement_match/relevance should still evaluate normally as if $6.5M were the real figure.

USER REQUEST:
{user_prompt}

UPLOADED FILES:
{uploaded_files}

SOURCE TEXT (extracted content from the uploaded files, used to verify factual claims - if it says "(none provided)", skip factual_accuracy and mark it not_checked):
{source_text}

AI OUTPUT:
{ai_output}

STRUCTURAL CHECKS (pass/fail):
1. requirement_match - did the output satisfy explicit counts and named instructions (e.g. "10 slides" means count them)? Ignore content depth and number correctness here.
2. missing_information - are all specifically requested topics/sections present at all? Ignore counts, formatting, and number correctness here.
3. completeness - does the output make use of every uploaded file (referencing its name or clearly drawing on its content counts as "used")? Ignore whether the figures used are correct here.
4. formatting - does the output match the requested STRUCTURE (e.g. slide-by-slide vs a paragraph)? Ignore content depth and number correctness here.

QUALITY CHECKS (score 1-5; pass is true only if score is 3 or higher):
5. fluency - is the writing grammatically correct and natural to read?
6. coherence - does the output flow logically without self-contradiction?
7. relevance - does the content stay on-topic for what was asked (subject matter, not number correctness), without irrelevant padding?
8. factual_accuracy - do specific claims and numbers in the AI OUTPUT match the SOURCE TEXT? This is high-severity: any invented figure not supported by SOURCE TEXT should fail it. If SOURCE TEXT is "(none provided)", set "pass": true and "reason": "not_checked - no source text provided".

INFORMATIONAL ONLY (does not affect pass/fail):
9. creativity - how original or engaging is the framing, scored 1-5?

Respond with ONLY valid JSON, no other text, in exactly this shape:
{{
  "requirement_match": {{"pass": true or false, "reason": "one short sentence"}},
  "missing_information": {{"pass": true or false, "reason": "one short sentence"}},
  "completeness": {{"pass": true or false, "reason": "one short sentence"}},
  "formatting": {{"pass": true or false, "reason": "one short sentence"}},
  "fluency": {{"score": 1, "pass": true or false, "reason": "one short sentence"}},
  "coherence": {{"score": 1, "pass": true or false, "reason": "one short sentence"}},
  "relevance": {{"score": 1, "pass": true or false, "reason": "one short sentence"}},
  "factual_accuracy": {{"pass": true or false, "reason": "one short sentence"}},
  "creativity": {{"score": 1, "reason": "one short sentence"}},
  "feedback": "combined, actionable instructions for whatever failed, empty string if everything passed"
}}
"""

GATING_CHECKS = [
    "requirement_match", "missing_information", "completeness", "formatting",
    "fluency", "coherence", "relevance", "factual_accuracy",
]

# Expected temperature ranges by task style. This is a plain sanity check, not
# an LLM judgment - temperature is a generation parameter, not something
# readable from the finished text, so we compare it directly as a number.
TEMPERATURE_RANGES = {
    "creative": (0.7, 1.0),
    "balanced": (0.4, 0.7),
    "deterministic": (0.0, 0.3),
}


def check_temperature(task_style, temperature_used):
    """Fast, free, non-LLM check: was the temperature actually used at
    generation time appropriate for this task's style? Returns None if either
    argument is missing (i.e. the caller isn't using this check)."""
    if task_style is None or temperature_used is None:
        return None
    if task_style not in TEMPERATURE_RANGES:
        return {"pass": True, "reason": "unknown task_style '" + str(task_style) + "', skipping check"}

    low, high = TEMPERATURE_RANGES[task_style]
    if low <= temperature_used <= high:
        return {
            "pass": True,
            "reason": "temperature " + str(temperature_used) + " is within the expected "
                      + task_style + " range (" + str(low) + "-" + str(high) + ")",
        }
    return {
        "pass": False,
        "reason": "temperature " + str(temperature_used) + " is outside the expected "
                  + task_style + " range (" + str(low) + "-" + str(high) + ")",
        "suggested_temperature": round((low + high) / 2, 2),
    }


def guardian_check(user_prompt, uploaded_files, ai_output, source_text=None,
                    task_style=None, temperature_used=None):
    """Run the Granite Guardian validation pass on a generated output. Returns a dict.

    source_text is optional: pass the actual extracted text from the uploaded
    files (not just filenames) to enable the factual_accuracy check. Without
    it, factual_accuracy is skipped automatically and always passes.

    task_style + temperature_used are optional: pass both (e.g. "creative"/0.85
    or "deterministic"/0.1) to enable a free, instant temperature sanity check
    alongside the LLM-judged checks. Without them, this check is skipped.
    """
    prompt = GUARDIAN_PROMPT_TEMPLATE.format(
        user_prompt=user_prompt,
        uploaded_files=", ".join(uploaded_files) if uploaded_files else "None",
        source_text=source_text if source_text else "(none provided)",
        ai_output=ai_output,
    )

    try:
        raw = call_granite(prompt)
    except GuardianConnectionError as e:
        return {
            "pass": False,
            "failed_checks": ["guardian_unavailable"],
            "feedback": str(e),
            "details": {},
            "creativity": None,
        }

    # Granite sometimes wraps JSON in markdown fences - strip those if present
    cleaned = raw.strip()
    if cleaned.startswith("```"):
        cleaned = cleaned.strip("`")
        if cleaned.startswith("json"):
            cleaned = cleaned[4:].strip()

    try:
        parsed = json.loads(cleaned)
    except json.JSONDecodeError:
        return {
            "pass": False,
            "failed_checks": ["parse_error"],
            "feedback": "Guardian response wasn't valid JSON: " + raw,
        }

    failed = [name for name in GATING_CHECKS if not parsed.get(name, {}).get("pass", False)]
    feedback = parsed.get("feedback", "")
    details = {name: parsed.get(name) for name in GATING_CHECKS}

    temp_result = check_temperature(task_style, temperature_used)
    if temp_result is not None:
        details["temperature"] = temp_result
        if not temp_result["pass"]:
            failed.append("temperature")
            note = ("Regenerate with temperature set closer to " + str(temp_result["suggested_temperature"])
                    + " - note this requires changing the API call's temperature parameter, "
                      "not just adding text to the prompt.")
            feedback = (feedback + " " + note).strip()

    return {
        "pass": len(failed) == 0,
        "failed_checks": failed,
        "feedback": feedback,
        "details": details,
        "creativity": parsed.get("creativity"),
    }


def guardian_check_from_orchestrator(orchestrator_output, uploaded_files=None, source_text=None):
    """Adapter for the Orchestrator's actual output shape:
    {
        "optimized_prompt": "...",
        "selected_model": "granite4.1:8b",
        "task_style": "creative",
        "temperature_used": 0.8,
        "generated_output": {
            "type": "text" or "file",
            "files": ["deck.pptx"],   # populated when actual files were produced
            "text": "..."            # populated when there's text content/summary
        }
    }

    uploaded_files and source_text come from whoever collects the user's
    input / Task 2 - pass them in here if available. Missing either one is
    fine: completeness and factual_accuracy just have less to check against.

    generated_output["text"] (when present) is used as a fallback. But when
    generated_output["files"] has real files, Guardian now extracts and
    judges their ACTUAL content instead (pptx/docx/pdf/xlsx/txt/csv/md) -
    that's the real deliverable, so it takes priority over any summary text.
    Formats without an extractor yet are flagged as "unvalidated" rather
    than silently failed, since that's a gap in this tool, not necessarily
    a problem with the output itself.

    Returns a single JSON-ready dict for Task 3 to consume directly:
    {
        "pass": true/false,
        "failed_checks": [...],
        "feedback": "human-readable explanation",
        "retry_prompt": "ready to resend as-is" (or null if pass is true),
        "retry_temperature": 0.85 (or null if pass is true)
    }
    She only needs to check "pass" and, if false, resend retry_prompt with
    retry_temperature - no interpretation logic needed on her side.
    """
    missing = [k for k in ("optimized_prompt", "generated_output") if k not in orchestrator_output]
    if missing:
        return {
            "pass": False,
            "failed_checks": ["malformed_orchestrator_output"],
            "feedback": "Orchestrator output is missing required field(s): " + ", ".join(missing),
            "retry_prompt": None,
            "retry_temperature": None,
        }

    user_prompt = orchestrator_output["optimized_prompt"]
    temperature_used = orchestrator_output.get("temperature_used")
    task_style = orchestrator_output.get("task_style")
    generated = orchestrator_output["generated_output"]

    # generated_output can be the old plain-string shape, or the new
    # {"type", "files", "text"} object - handle both so this doesn't break
    # again if the shape shifts slightly.
    if isinstance(generated, dict):
        output_type = generated.get("type", "text")
        output_text = generated.get("text") or ""
        output_files = generated.get("files") or []
    else:
        output_type = "text"
        output_text = generated or ""
        output_files = []

    # --- Read and validate real files, if any were produced -------------
    file_reports = []
    extracted_texts = []
    broken_files = []
    unsupported_formats = []

    if output_type == "file" and output_files:
        for path in output_files:
            extraction = extract_file_content(path)
            file_reports.append({"file": path, "kind": extraction["kind"], "error": extraction["error"]})

            if extraction["kind"] == "text":
                extracted_texts.append("--- " + os.path.basename(path) + " ---\n" + extraction["content"])
            elif extraction["kind"] == "unsupported":
                unsupported_formats.append(path)
            elif extraction["kind"] == "error":
                broken_files.append(path + " (" + str(extraction["error"]) + ")")

    # Real extracted file content is the ground truth of what the user
    # actually receives, so it takes priority over any summary text.
    content_to_judge = "\n\n".join(extracted_texts) if extracted_texts else output_text

    if content_to_judge.strip():
        result = guardian_check(
            user_prompt=user_prompt,
            uploaded_files=uploaded_files or [],
            ai_output=content_to_judge,
            source_text=source_text,
            task_style=task_style,
            temperature_used=temperature_used,
        )
    else:
        # Nothing textual to hand the LLM checks - skip them gracefully
        # rather than judging an empty string.
        result = {"pass": True, "failed_checks": [], "feedback": "", "details": {}, "creativity": None}

    # --- Fast, non-LLM checks about file delivery itself -----------------
    if output_type == "file":
        if not output_files:
            result["pass"] = False
            result["failed_checks"] = result["failed_checks"] + ["output_delivered"]
            result["feedback"] = (
                result["feedback"]
                + " Output type was 'file' but no files were actually produced - regenerate and ensure a file is attached."
            ).strip()
        else:
            result["file_reports"] = file_reports

            if broken_files:
                result["pass"] = False
                result["failed_checks"] = result["failed_checks"] + ["file_unreadable"]
                result["feedback"] = (
                    result["feedback"] + " Could not read: " + "; ".join(broken_files)
                ).strip()

            if unsupported_formats:
                # A gap in this tool, not necessarily a bad output - don't
                # fail for it, but make it visible rather than silent.
                result["unvalidated_files"] = unsupported_formats
                result["feedback"] = (
                    result["feedback"]
                    + " Note: no content validator yet for: " + ", ".join(unsupported_formats)
                    + " (not held against pass/fail)."
                ).strip()

    retry_prompt, retry_temperature = build_retry_payload(user_prompt, temperature_used, result)
    result["retry_prompt"] = retry_prompt
    result["retry_temperature"] = retry_temperature

    return result


# ---- Mock test cases so you can test WITHOUT the rest of the pipeline built yet ----

TEST_CASES = [
    {
        "name": "Good output - should PASS",
        "user_prompt": "Create a 10-slide sales presentation covering revenue, expenses, and profit.",
        "uploaded_files": ["Q3_Sales_Report.pdf"],
        "ai_output": (
            "Slide 1: Title - Q3 Sales Performance. "
            "Slide 2: Revenue Overview - Per Q3_Sales_Report.pdf, revenue reached $4.2M, up 12% QoQ. "
            "Slide 3: Revenue Trends - Monthly growth accelerated from 8% in July to 15% in September, per the report. "
            "Slide 4: Expenses Breakdown - Q3_Sales_Report.pdf shows total expenses of $2.1M, an 8% increase. "
            "Slide 5: Cost Analysis - Marketing and logistics costs drove most of the expense growth. "
            "Slide 6: Profit Summary - Net profit was $2.1M, a 16% margin, based on report figures. "
            "Slide 7: Profit Trends - Profit margin improved steadily each month. "
            "Slide 8: Risks - Rising logistics costs pose a risk to margins. "
            "Slide 9: Recommendations - Negotiate better logistics rates and expand high-margin product lines. "
            "Slide 10: Conclusion - Strong Q3 performance driven by revenue growth outpacing expense growth."
        ),
    },
    {
        "name": "Too few slides - should FAIL (Requirement Match) ONLY",
        "user_prompt": "Create a 10-slide sales presentation covering revenue, expenses, and profit.",
        "uploaded_files": ["Q3_Sales_Report.pdf"],
        "ai_output": (
            "Slide 1: Revenue Overview - Per Q3_Sales_Report.pdf, revenue reached $4.2M, up 12% QoQ. "
            "Slide 2: Expenses Breakdown - Q3_Sales_Report.pdf shows total expenses of $2.1M, an 8% increase. "
            "Slide 3: Profit Summary - Net profit was $2.1M, a 16% margin, based on report figures. "
            "Slide 4: Conclusion - Strong Q3 performance overall."
        ),
    },
    {
        "name": "Missing a required topic - should FAIL (Missing Information) ONLY",
        "user_prompt": "Create a presentation covering revenue, expenses, and profit.",
        "uploaded_files": ["Q3_Sales_Report.pdf"],
        "ai_output": (
            "Slide 1: Revenue Overview - Per Q3_Sales_Report.pdf, revenue reached $4.2M, up 12% QoQ. "
            "Slide 2: Revenue Trends - Monthly growth accelerated from 8% in July to 15% in September, per the report. "
            "Slide 3: Revenue by Region - North America led growth at 18%, per Q3_Sales_Report.pdf. "
            "Slide 4: Conclusion - Revenue performance was strong this quarter."
        ),
    },
    {
        "name": "Ignored an uploaded file - should FAIL (Completeness) ONLY",
        "user_prompt": "Create a presentation covering our sales and marketing performance using both reports.",
        "uploaded_files": ["Sales_Report.pdf", "Marketing_Report.pdf"],
        "ai_output": (
            "Slide 1: Sales Overview - Per Sales_Report.pdf, revenue reached $4.2M, up 12% QoQ. "
            "Slide 2: Sales Trends - Monthly growth accelerated through the quarter, per Sales_Report.pdf. "
            "Slide 3: Expense Summary - Sales_Report.pdf shows costs of $2.1M for the period. "
            "Slide 4: Conclusion - Strong sales performance driven by regional growth."
        ),
    },
    {
        "name": "Wrong format - should FAIL (Formatting) ONLY",
        "user_prompt": "Create a professional slide presentation covering revenue, expenses, and profit using the uploaded report.",
        "uploaded_files": ["Report.pdf"],
        "ai_output": (
            "Here is a summary: According to Report.pdf, revenue grew 10% this quarter to $3.8M. "
            "Expenses rose slightly to $2.0M, and net profit came in at $1.8M, a healthy margin. "
            "Overall this was a strong quarter for the business."
        ),
    },
    {
        "name": "Fabricated numbers - should FAIL (factual_accuracy) ONLY",
        "user_prompt": "Summarize the Q3 sales report.",
        "uploaded_files": ["Q3_Sales_Report.pdf"],
        "source_text": (
            "Q3 Sales Report: Total revenue was $4.2M, up 12% year over year. "
            "Total expenses were $2.8M. Net profit was $1.4M with a 33% margin."
        ),
        "ai_output": (
            "Based on the Q3 Sales Report, Q3 revenue reached $6.5M, a 40% increase year over year, "
            "with expenses of only $1.2M and a record profit margin of 55%."
        ),
    },
    {
        "name": "Incoherent writing - should FAIL (fluency / coherence) ONLY",
        "user_prompt": "Summarize the Q3 sales report in a short paragraph.",
        "uploaded_files": ["Q3_Sales_Report.pdf"],
        "source_text": (
            "Q3 Sales Report: Total revenue was $4.2M, up 12% year over year. "
            "Total expenses were $2.8M. Net profit was $1.4M with a 33% margin."
        ),
        "ai_output": (
            "revenue up money good quarter also expenses but the the profit was is happen "
            "good margin number thing increase yes report say so."
        ),
    },
    {
        "name": "Creative task run too cold - should FAIL (temperature) ONLY",
        "user_prompt": "Write a fun, imaginative product tagline for our new eco-friendly water bottle.",
        "uploaded_files": [],
        "task_style": "creative",
        "temperature_used": 0.1,
        "ai_output": "Sip the future: hydration reimagined for a thirstier planet.",
    },
    {
        "name": "Deterministic task run too hot - should FAIL (temperature) ONLY",
        "user_prompt": "Extract the exact total revenue figure from the uploaded report.",
        "uploaded_files": ["Q3_Sales_Report.pdf"],
        "source_text": (
            "Q3 Sales Report: Total revenue was $4.2M, up 12% year over year. "
            "Total expenses were $2.8M. Net profit was $1.4M with a 33% margin."
        ),
        "task_style": "deterministic",
        "temperature_used": 0.9,
        "ai_output": "The total revenue was $4.2M.",
    },
]


def build_retry_payload(user_prompt, temperature_used, guardian_result):
    """Given a FAILED guardian_result, build the exact retry prompt and
    temperature to regenerate with. Returns (None, None) if guardian_result
    already passed - nothing to retry."""
    if guardian_result["pass"]:
        return None, None

    retry_temperature = temperature_used
    temp_check = guardian_result.get("details", {}).get("temperature")
    if temp_check and not temp_check.get("pass", True) and "suggested_temperature" in temp_check:
        retry_temperature = temp_check["suggested_temperature"]

    retry_prompt = (
        user_prompt
        + "\n\nThe previous attempt had issues. Please fix the following and regenerate:\n"
        + guardian_result["feedback"]
    )
    return retry_prompt, retry_temperature


# ============================================================================
# REAL PIPELINE FILE I/O
# ----------------------------------------------------------------------------
# Everything needed to run Guardian against the actual handoff files:
# Payal's generated_result.json (Task 3's output) and Shalmalee's
# cleaned_input.json (Task 2's extracted document). Writes timestamped result
# files instead of just printing, since a real pipeline run needs a record on
# disk - both for the frontend to pick up (final_output.json) and for Task 3
# to pick up on retry (guardian_retry.json), since Task 3 has no memory
# between calls and needs the full prompt handed back to it explicitly.
# ============================================================================

def load_json_file(path):
    """Load a JSON file from disk with a clear, actionable error message if
    it's missing or malformed, rather than letting a raw exception crash the
    script or silently return garbage."""
    if not os.path.exists(path):
        raise GuardianConnectionError("Could not find file: " + path)
    try:
        with open(path, "r", encoding="utf-8") as f:
            return json.load(f)
    except json.JSONDecodeError as e:
        raise GuardianConnectionError("File '" + path + "' is not valid JSON: " + str(e))


# Large source documents (e.g. a 324-page PDF at 500k+ characters) are far
# more than a local 8B model's context window can use productively. Truncate
# to a safe budget and say so plainly in the text itself, rather than
# silently sending a wall of text that gets cut off mid-token by Ollama.
MAX_SOURCE_TEXT_CHARS = 12000


def _truncate_source_text(text, max_chars=MAX_SOURCE_TEXT_CHARS):
    if not text:
        return ""
    if len(text) <= max_chars:
        return text
    return (
        text[:max_chars]
        + "\n\n[... truncated for length - " + str(len(text)) + " total characters in the original document ...]"
    )


def _timestamped_path(label):
    stamp = datetime.now().strftime("%Y%m%d_%H%M%S_%f")
    return os.path.join(RESULTS_DIR, label + "_" + stamp + ".json")


def _append_run_log(entry):
    with open(RUN_LOG_PATH, "a", encoding="utf-8") as f:
        f.write(json.dumps(entry) + "\n")


def save_guardian_result(result, orchestrator_output, output_path=None):
    """Save the result of a real pipeline run (--files mode). Always writes
    THREE files with FIXED names, overwritten every run, so the orchestrator
    can read a deterministic path instead of having to find "the latest" one:

      - generated_output.json - the full Guardian verdict (all 9 checks,
        details, feedback) - for debugging.
      - final.json - {files, text} for the frontend. Populated with the real
        generated content when Guardian PASSED; empty ({"files": [],
        "text": ""}) when it FAILED, since there's nothing valid to hand the
        frontend yet.
      - retry.json - {retry_prompt, retry_temperature, failed_checks,
        feedback} for Task 3 to regenerate with. Populated when Guardian
        FAILED; empty/null when it PASSED, since there's nothing to retry.

    History isn't lost by overwriting these - every run still appends a line
    to RUN_LOG_PATH (guardian_results/run_log.jsonl) regardless, so you can
    trace what happened on past runs even though only the latest result
    lives at these three fixed paths.

    orchestrator_output is Payal's raw generated_result.json content - used
    to pull the actual generated files/text for final.json (note: those live
    nested under generated_output, not at the top level).

    Returns the dict of paths written.
    """
    generated_output = orchestrator_output.get("generated_output", {})
    if not isinstance(generated_output, dict):
        generated_output = {}
    generated_files = generated_output.get("files", [])
    generated_text = generated_output.get("text", "")

    passed = bool(result.get("pass"))

    generated_output_path = output_path or os.path.join(RESULTS_DIR, "generated_output.json")
    with open(generated_output_path, "w", encoding="utf-8") as f:
        json.dump(result, f, indent=2)

    final_path = os.path.join(RESULTS_DIR, "final.json")
    final_output = {
        "files": generated_files if passed else [],
        "text": generated_text if passed else "",
    }
    with open(final_path, "w", encoding="utf-8") as f:
        json.dump(final_output, f, indent=2)

    retry_path = os.path.join(RESULTS_DIR, "retry.json")
    retry_output = {
        "retry_prompt": result.get("retry_prompt"),           # already None when passed
        "retry_temperature": result.get("retry_temperature"),  # already None when passed
        "failed_checks": result.get("failed_checks") if not passed else [],
        "feedback": result.get("feedback") if not passed else "",
    }
    with open(retry_path, "w", encoding="utf-8") as f:
        json.dump(retry_output, f, indent=2)

    written = {
        "generated_output": generated_output_path,
        "final_output": final_path,
        "retry_output": retry_path,
    }

    _append_run_log({
        "timestamp": datetime.now().isoformat(timespec="seconds"),
        "mode": "files",
        "pass": passed,
        "failed_checks": result.get("failed_checks"),
        **written,
    })

    return written


def save_run_results(all_results, output_path=None, label="guardian_run"):
    """Save results from the mock-test-case / demo run (not the real-files
    mode - that's save_guardian_result above). Just dumps the whole results
    dict to a timestamped file plus a run-log line; there's no single
    "generated file" to extract a final_output from when this is a batch of
    unrelated test cases."""
    path = output_path or _timestamped_path(label)
    with open(path, "w", encoding="utf-8") as f:
        json.dump(all_results, f, indent=2)

    _append_run_log({
        "timestamp": datetime.now().isoformat(timespec="seconds"),
        "mode": "test_run",
        "output": path,
    })
    return path


def guardian_check_from_files(generated_result_path="generated_result.json",
                               cleaned_input_path="cleaned_input.json"):
    """End-to-end entry point for a real pipeline run.

    - generated_result_path: Task 3's (Payal's) output - optimized_prompt /
      selected_model / task_style / temperature_used / generated_output.
    - cleaned_input_path: Task 2's (Shalmalee's) cleaned, extracted document -
      original_prompt / document_text / requirements / metadata (including
      the list of source filenames under metadata.documents).

    Loads both files, runs guardian_check_from_orchestrator, then hands the
    result to save_guardian_result() to write (fixed filenames, overwritten
    each run - see save_guardian_result docstring for details):
      - generated_output.json - always, the full verdict
      - retry.json            - populated if it FAILED, empty/null if it
        PASSED - everything Task 3 needs to regenerate (prompt with feedback
        folded in, and temperature), since Task 3 has no memory between calls
      - final.json             - populated if it PASSED, empty if it FAILED -
        the clean {files, text} payload that actually goes to the frontend

    Returns (result, written_paths).
    """
    orchestrator_output = load_json_file(generated_result_path)
    cleaned_input = load_json_file(cleaned_input_path)

    document_text = cleaned_input.get("document_text", "")
    source_text = _truncate_source_text(document_text)

    documents = cleaned_input.get("metadata", {}).get("documents", [])
    uploaded_files = [d.get("filename", "unknown") for d in documents] or None

    result = guardian_check_from_orchestrator(
        orchestrator_output,
        uploaded_files=uploaded_files,
        source_text=source_text if source_text.strip() else None,
    )

    written = save_guardian_result(result, orchestrator_output)

    if result["pass"]:
        print("Guardian PASSED.")
    else:
        print("Guardian FAILED (" + ", ".join(result["failed_checks"]) + ").")
    print("generated_output -> " + written["generated_output"])
    print("final_output     -> " + written["final_output"])
    print("retry_output     -> " + written["retry_output"])

    return result, written


def run_with_guardian(user_prompt, uploaded_files, generate_fn, source_text=None,
                       task_style=None, temperature_used=None):
    """Orchestrates the full flow: generate -> guardian_check -> if it fails,
    regenerate ONCE with the feedback folded in -> guardian_check again -> return.

    generate_fn: a function with signature (prompt, temperature) -> str.
    This stands in for "call the chosen model" - useful for testing your own
    Guardian in isolation. In the real pipeline, Task 3 owns the retry loop
    herself and calls guardian_check_from_orchestrator directly instead - see
    that function below.
    """
    output = generate_fn(user_prompt, temperature_used)
    result = guardian_check(user_prompt, uploaded_files, output, source_text, task_style, temperature_used)

    if not result["pass"]:
        retry_prompt, retry_temperature = build_retry_payload(user_prompt, temperature_used, result)
        output = generate_fn(retry_prompt, retry_temperature)
        result = guardian_check(user_prompt, uploaded_files, output, source_text, task_style, retry_temperature)
        result["retried"] = True
    else:
        result["retried"] = False

    result["final_output"] = output
    return result


def _demo_generate_fn(prompt, temperature):
    """Toy stand-in for a real model call, used only to demo the retry flow end
    to end. Returns a deliberately thin answer on the first call, and a full
    one once it sees this is a retry - mimicking a model fixing itself after
    feedback. Replace with Task 3's real generation function when it's ready."""
    if "previous attempt had issues" in prompt:
        return (
            "Slide 1: Title - Q3 Sales Performance. "
            "Slide 2: Revenue overview - $4.2M, up 12% YoY, per Q3_Sales_Report.pdf. "
            "Slide 3: Revenue trends - steady growth across regions. "
            "Slide 4: Expenses breakdown - $2.8M total expenses. "
            "Slide 5: Cost analysis - expenses grew slower than revenue. "
            "Slide 6: Profit summary - net profit of $1.4M, 33% margin. "
            "Slide 7: Profit trends - margin improved through the quarter. "
            "Slide 8: Risks - supply chain delays noted in the report. "
            "Slide 9: Recommendations - invest in logistics automation. "
            "Slide 10: Conclusion - a strong quarter overall."
        )
    return "Slide 1: Revenue. Slide 2: Expenses. Slide 3: Profit."


# Curated for live demos: one hallucination catch, one temperature catch -
# enough variety to show real breadth without running all 9 on stage.
DEMO_CASE_NAMES = [
    "Fabricated numbers - should FAIL (factual_accuracy) ONLY",
    "Creative task run too cold - should FAIL (temperature) ONLY",
]

if __name__ == "__main__":
    # --files mode: run against the REAL pipeline handoff files instead of
    # the mock test cases below. Usage:
    #   python granite_guardian.py --files
    #   python granite_guardian.py --files path/to/generated_result.json path/to/cleaned_input.json
    if "--files" in sys.argv:
        file_args = [a for a in sys.argv[1:] if a != "--files"]
        gen_path = file_args[0] if len(file_args) > 0 else "generated_result.json"
        clean_path = file_args[1] if len(file_args) > 1 else "cleaned_input.json"

        if not check_ollama_status():
            print("\nFix the issue above, then rerun this script.")
            raise SystemExit(1)
        warm_up_model()

        try:
            guardian_check_from_files(gen_path, clean_path)
        except GuardianConnectionError as e:
            print("Could not run against files: " + str(e))
            raise SystemExit(1)

        raise SystemExit(0)

    if not check_ollama_status():
        print("\nFix the issue above, then rerun this script.")
        raise SystemExit(1)

    warm_up_model()

    demo_mode = "--demo" in sys.argv
    cases_to_run = [c for c in TEST_CASES if c["name"] in DEMO_CASE_NAMES] if demo_mode else TEST_CASES

    all_results = {}
    for case in cases_to_run:
        print("\n=== " + case["name"] + " ===")
        result = guardian_check(
            case["user_prompt"],
            case["uploaded_files"],
            case["ai_output"],
            case.get("source_text"),
            case.get("task_style"),
            case.get("temperature_used"),
        )
        print(json.dumps(result, indent=2))
        all_results[case["name"]] = result

    print("\n\n=== DEMO: full generate -> guardian -> retry flow ===")
    demo_result = run_with_guardian(
        "Create a 10-slide sales presentation covering revenue, expenses, and profit.",
        ["Q3_Sales_Report.pdf"],
        _demo_generate_fn,
        source_text=(
            "Q3 Sales Report: Total revenue was $4.2M, up 12% year over year. "
            "Total expenses were $2.8M. Net profit was $1.4M with a 33% margin."
        ),
    )
    print(json.dumps(demo_result, indent=2))
    all_results["DEMO: full generate -> guardian -> retry flow"] = demo_result

    if not demo_mode:
        print("\n\n=== ORCHESTRATOR ADAPTER: exact shape from her screenshot ===")
        her_example_output = {
            "optimized_prompt": "Summarize the provided document, focusing on how AI is used across industries.",
            "selected_model": "granite3.1-moe:3b",
            "task_style": "creative",
            "temperature_used": 0.9,
            "generated_output": {
                "type": "file",
                "files": [],  # empty, exactly as in her screenshot
                "text": "- Healthcare Sector:\n  - AI enhances diagnostics through faster image analysis.",
            },
        }
        result_a = guardian_check_from_orchestrator(her_example_output)
        print(json.dumps(result_a, indent=2))
        all_results["ORCHESTRATOR ADAPTER: exact shape from her screenshot"] = result_a

        print("\n\n=== ORCHESTRATOR ADAPTER: file type WITH a file actually attached ===")
        her_example_output_fixed = dict(her_example_output)
        her_example_output_fixed["generated_output"] = dict(her_example_output["generated_output"])
        her_example_output_fixed["generated_output"]["files"] = ["ai_industries_summary.pptx"]
        result_b = guardian_check_from_orchestrator(her_example_output_fixed)
        print(json.dumps(result_b, indent=2))
        all_results["ORCHESTRATOR ADAPTER: file type WITH a file actually attached"] = result_b

    log_path = save_run_results(all_results, label="guardian_demo_run" if demo_mode else "guardian_test_run")
    print("\n\nAll results for this run saved to " + log_path)